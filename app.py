import streamlit as st
import openai
import tempfile
import os
from PyPDF2 import PdfReader
from docx import Document
from docx import Document as DocxWriter
import requests
from bs4 import BeautifulSoup
import re # For regex to highlight figures
from googlesearch import search # This is the correct way to import the provided tool
import difflib # For comparing text differences

# Libraries for Excel and PowerPoint - ensure these are installed via pip
try:
    import openpyxl
except ImportError:
    st.warning("`openpyxl` not found. Please install it: `pip install openpyxl`")
try:
    from pptx import Presentation
except ImportError:
    st.warning("`python-pptx` not found. Please install it: `pip install python-pptx`")


# --- Text Extraction Functions ---
def extract_pdf_text(file):
    """Extracts text from a PDF file."""
    try:
        reader = PdfReader(file)
        return "".join(page.extract_text() or "" for page in reader.pages)[:8000]
    except Exception as e:
        return f"[Error extracting PDF text: {e}]"

def extract_docx_text(file):
    """Extracts text from a DOCX file."""
    try:
        doc = Document(file)
        return "\n".join(para.text for para in doc.paragraphs)[:8000]
    except Exception as e:
        return f"[Error extracting DOCX text: {e}]"

def extract_excel_text(file):
    """Extracts text from an XLSX file (all sheets)."""
    try:
        workbook = openpyxl.load_workbook(file)
        full_text = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            full_text.append(f"--- Sheet: {sheet_name} ---")
            for row in sheet.iter_rows():
                row_values = [str(cell.value) if cell.value is not None else "" for cell in row]
                full_text.append("\t".join(row_values))
        return "\n".join(full_text)[:8000]
    except Exception as e:
        return f"[Error extracting XLSX text: {e}]"

def extract_ppt_text(file):
    """Extracts text from a PPTX file (all slides)."""
    try:
        prs = Presentation(file)
        full_text = []
        for i, slide in enumerate(prs.slides):
            full_text.append(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return "\n".join(full_text)[:8000]
    except Exception as e:
        return f"[Error extracting PPTX text: {e}]"

# --- Web Scraping Functions (Internal, not exposed in UI) ---
def fetch_text_from_url(url):
    """Scrapes paragraphs from a given URL."""
    try:
        response = requests.get(url)
        soup = BeautifulSoup(response.text, "html.parser")
        paragraphs = soup.find_all('p')
        return "\n".join(p.get_text(strip=True) for p in paragraphs[:10])
    except Exception as e:
        return f"[Error fetching URL content: {e}]"

def fetch_legal_examples(query_term):
    """
    Fetches legal examples from various online sources using Google Search.
    This function is internal and its sources are not mentioned in the UI.
    """
    all_snippets = []

    # Search for EDGAR filings
    edgar_query = f"site:sec.gov/Archives/edgar/data {query_term} agreement"
    edgar_results = list(search(edgar_query, num_results=5))
    if edgar_results:
        all_snippets.extend([f"Edgar filing content related to {query_term}" for _ in edgar_results[:3]])

    # Search for Law Firm website content
    lawfirm_query = f"site:.com law firm {query_term} contract clauses OR template"
    lawfirm_results = list(search(lawfirm_query, num_results=5))
    if lawfirm_results:
        all_snippets.extend([f"Law firm content related to {query_term}" for _ in lawfirm_results[:3]])

    # Fetch from LawInsider (direct scrape as it's a known structure)
    try:
        lawinsider_url = "https://www.lawinsider.com/clause/scope-of-work"
        response = requests.get(lawinsider_url)
        soup = BeautifulSoup(response.text, "html.parser")
        clauses = soup.select(".clause-body")
        all_snippets.extend([clause.get_text(strip=True) for clause in clauses[:5]])
    except Exception as e:
        print(f"Error fetching from LawInsider: {e}")

    # Search wider internet for similar agreements
    general_query = f"'{query_term}' agreement examples OR template OR clauses"
    general_results = list(search(general_query, num_results=5))
    if general_results:
        all_snippets.extend([f"General agreement content for {query_term}" for _ in general_results[:3]])

    return "\n---\n".join(all_snippets[:10]) # Limit to top 10 snippets for prompt size

# --- Text Comparison Function ---
def highlight_differences(original_text, refined_text):
    """
    Compares original and refined text and returns HTML with highlighted changes.
    Also cleans up figure highlighting to show only the values.
    """
    # Clean up figure highlighting in both texts before comparison
    def clean_figures(text):
        # Remove [FIGURE: ] wrapper and keep only the value
        return re.sub(r"<span style='background-color: yellow; padding: 2px 4px; border-radius: 3px;'>\[FIGURE: (.*?)\]</span>", r'<span style="background-color: yellow; padding: 2px 4px; border-radius: 3px;">\1</span>', text)
    
    original_text = clean_figures(original_text)
    refined_text = clean_figures(refined_text)
    
    original_lines = original_text.split('\n')
    refined_lines = refined_text.split('\n')
    
    highlighted_refined = []
    
    for line in difflib.ndiff(original_lines, refined_lines):
        if line.startswith('  '):  # unchanged line
            highlighted_refined.append(line[2:])
        elif line.startswith('+ '):  # added/modified line
            highlighted_line = f"<div style='background-color: #e8f5e8; border-left: 4px solid #4CAF50; padding: 8px; margin: 2px 0;'><strong>REFINED:</strong> {line[2:]}</div>"
            highlighted_refined.append(highlighted_line)
        elif line.startswith('- '):  # removed line (original)
            highlighted_line = f"<div style='background-color: #fff2f2; border-left: 4px solid #f44336; padding: 8px; margin: 2px 0; text-decoration: line-through;'><strong>ORIGINAL:</strong> {line[2:]}</div>"
            highlighted_refined.append(highlighted_line)
    
    return '\n'.join(highlighted_refined)

# --- LLM Interaction Function ---
def generate_sow(base_text, user_desc, role_preference, combined_examples, existing_sow=None, feedback=None, additional_context=""):
    """
    Generates or refines a Scope of Work (SoW) using an LLM.
    Adjusts content based on role preference (pro-vendor/pro-client).
    """
    # Define roles for the LLM prompt
    if role_preference == "Company as Service Provider (Pro-Vendor)":
        company_role = "Service Provider"
        client_role = "Company"
        dependencies_guidance = "The dependencies to be provided by the Company (client) should be extensive and clearly defined. Obligations of the Service Provider (your company) should be more generic and high-level."
    else: # Company as Service Recipient (Pro-Client)
        company_role = "Company"
        client_role = "Service Provider" # In this case, the 'Service Provider' is the external entity
        dependencies_guidance = "The dependencies to be provided by the Company (client) should be basic and minimal. Obligations of the Service Provider (the external entity) should be detailed and specific."

    examples_text = "\n---\n".join(combined_examples) if combined_examples else "None included"

    # Base prompt structure for initial generation or refinement
    base_prompt_template = f"""
You are a legal AI assistant specializing in contract drafting. Your task is to generate a comprehensive Scope of Work (SoW) document.

The Company in this SoW refers to the entity that is the {'service provider' if role_preference == 'Company as Service Provider (Pro-Vendor)' else 'service recipient'}.
The Service Provider in this SoW refers to the external entity that is the {'service recipient' if role_preference == 'Company as Service Provider (Pro-Vendor)' else 'service provider'}.

{dependencies_guidance}

Highlight all figures (numbers, percentages, currency amounts) by wrapping them in <span style='background-color: yellow; padding: 2px 4px; border-radius: 3px;'>[FIGURE: X]</span> tags to indicate they require independent validation. For example, a price of $10,000 should be <span style='background-color: yellow; padding: 2px 4px; border-radius: 3px;'>[FIGURE: $10,000]</span>.

The SoW should follow this exact structure:

1. Description – What is being supplied or done.
2. Function – The business purpose or outcome the goods/services serve.
3. Price – Pricing structure, billing frequency, and payment terms.
4. Dependencies for the service to be provided by the Company/client.
5. Milestones – Key deliverables with corresponding deadlines or phases.
6. Warranties – Any performance guarantees, service warranties, or coverage periods.
7. Service Levels (if applicable) – SLAs, KPIs, uptime, penalties, or escalation paths.
8. Others – Any additional relevant clauses not captured above (e.g. assumptions, subcontracting, ownership of deliverables).
"""

    if existing_sow and feedback:
        # Refinement prompt
        prompt = f"""
{base_prompt_template}

---
Existing Scope of Work:
{existing_sow}

---
User Feedback for Refinement:
{feedback}

---
Based on the existing Scope of Work and the user's feedback, provide an updated version of the SoW. Retain all parts that have not been highlighted or are not relevant to the changes requested in the feedback. Ensure the updated SoW still adheres to the specified structure and highlighting rules.
"""
    else:
        # Initial generation prompt
        prompt = f"""
{base_prompt_template}

---
User Description of Goods/Services and Business Context:
{user_desc}

---
Base Document Extract (from uploaded files):
{base_text}

---
User-Provided and Automatically Fetched Example SoW Clauses:
{examples_text}

---
Additional Context from Public Filings and Legal Resources:
{additional_context if additional_context else "No additional relevant context found."}

---
Generate a detailed Scope of Work (SoW) based on the provided information, adhering to the structure, role definitions, and highlighting rules. Also, suggest questions for missing or unclear details at the end of the generated SoW.
"""

    # Get OpenAI API key
    # Try st.secrets first (recommended for Streamlit Cloud), then os.getenv
    openai_api_key = st.secrets.get("OPENAI_API_KEY") or os.getenv("OPENAI_API_KEY")

    if not openai_api_key:
        st.error("OpenAI API key not found. Please set it in Streamlit secrets or as an environment variable.")
        st.stop()

    client = openai.OpenAI(base_url="https://generativelanguage.googleapis.com/v1beta/openai/", api_key=openai_api_key)

    response = client.chat.completions.create(
        model="gemini-2.5-flash",
        messages=[
            {"role": "system", "content": "You are a legal AI assistant specializing in contract drafting."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.7
    )
    return response.choices[0].message.content

# --- DOCX Export Function ---
def export_to_docx(content, file_name="Scope_of_Work.docx"):
    """Exports the given content to a DOCX file."""
    doc = DocxWriter()
    doc.add_heading("Generated Scope of Work", level=1)
    # Remove the HTML highlighting tags before exporting to DOCX
    clean_content = re.sub(r"<span style='background-color: yellow;[^>]*?>\[FIGURE: (.*?)\]</span>", r"\1", content)
    for paragraph in clean_content.split('\n'):
        doc.add_paragraph(paragraph)
    temp_path = tempfile.NamedTemporaryFile(delete=False, suffix=".docx")
    doc.save(temp_path.name)
    return temp_path.name

# --- Streamlit UI ---
st.title("AI Scope of Work (SoW) Generator")

# Document Upload Section
uploaded_file = st.file_uploader(
    "Upload all relevant/to-date client presentations, proposals, scope documents, and even base contracts (PDF, DOCX, XLSX, PPTX)",
    type=["pdf", "docx", "xlsx", "pptx"]
)
user_desc = st.text_area("Describe the goods/services and business context")

# Role Selection
role_preference = st.radio(
    "Generate SoW for your Company as:",
    ("Company as Service Provider (Pro-Vendor)", "Company as Service Recipient (Pro-Client)")
)

# Optional External Content sections
st.markdown("---")
st.subheader("Optional External Content")
custom_examples_input = st.text_area("Paste your own SoW clauses or content here (optional)")
external_url = st.text_input("Paste a URL to extract external SoW-style clauses (optional)")

# Keyword for additional search
st.markdown("---")
st.subheader("Additional Search Context")
search_keyword = st.text_input("Keyword to search", value=user_desc)


if st.button("Generate SoW"):
    if uploaded_file and user_desc:
        with st.spinner("Extracting text and gathering legal context..."):
            # Extract uploaded content
            file_extension = uploaded_file.name.split(".")[-1].lower()
            base_text = ""
            if file_extension == "pdf":
                base_text = extract_pdf_text(uploaded_file)
            elif file_extension == "docx":
                base_text = extract_docx_text(uploaded_file)
            elif file_extension == "xlsx":
                base_text = extract_excel_text(uploaded_file)
            elif file_extension == "pptx":
                base_text = extract_ppt_text(uploaded_file)
            else:
                st.error("Unsupported file type.")
                st.stop()

            # Collect all examples, including user-provided and URL-scraped
            combined_examples = []
            if custom_examples_input.strip():
                combined_examples.append(custom_examples_input.strip())
            if external_url.strip():
                fetched_url_content = fetch_text_from_url(external_url.strip())
                if fetched_url_content:
                    combined_examples.append(fetched_url_content)

            # Automatically comb through public documents and legal resources
            additional_context = fetch_legal_examples(search_keyword)

            # Generate SoW
            generated_sow_content = generate_sow(
                base_text,
                user_desc,
                role_preference,
                combined_examples=combined_examples,
                additional_context=additional_context
            )

            # Store the generated SoW in session state for refinement
            st.session_state.generated_sow = generated_sow_content
            st.session_state.base_text = base_text
            st.session_state.user_desc = user_desc
            st.session_state.role_preference = role_preference
            st.session_state.combined_examples = combined_examples
            st.session_state.additional_context = additional_context

            # Show result with highlighting
            st.subheader("Generated Scope of Work")
            st.markdown(generated_sow_content, unsafe_allow_html=True)

            # Offer download
            docx_path = export_to_docx(generated_sow_content, "Generated_Scope_of_Work.docx")
            with open(docx_path, "rb") as f:
                st.download_button("Download SoW as DOCX", f, file_name="Generated_Scope_of_Work.docx")
    else:
        st.warning("Please upload a document and provide a description to generate the SoW.")

# Iterative Refinement Section
st.markdown("---")
st.subheader("Refine SoW")

if 'generated_sow' in st.session_state and st.session_state.generated_sow:
    # Display the current SoW in a disabled text area
    st.text_area(
        "Current Scope of Work (scroll to view full content)",
        value=st.session_state.generated_sow,
        height=400,
        disabled=True,
        key="current_sow_display"
    )

    feedback_input = st.text_area(
        "Suggest a refinement (e.g., expand delivery details, adjust pricing terms, clarify dependencies):",
        key="feedback_for_refinement"
    )

    if st.button("Apply Refinement"):
        if feedback_input.strip():
            with st.spinner("Refining Scope of Work..."):
                # Store the original SoW before refinement
                original_sow = st.session_state.generated_sow
                
                refined_sow = generate_sow(
                    st.session_state.base_text,
                    st.session_state.user_desc,
                    st.session_state.role_preference,
                    combined_examples=st.session_state.combined_examples,
                    existing_sow=st.session_state.generated_sow,
                    feedback=feedback_input,
                    additional_context=st.session_state.additional_context
                )
                st.session_state.generated_sow = refined_sow # Update the stored SoW for next iteration

                st.success("Scope of Work updated based on your refinement. You can continue refining below.")
                
                # Show highlighted differences
                st.subheader("Refined Scope of Work with Changes Highlighted")
                highlighted_content = highlight_differences(original_sow, refined_sow)
                st.markdown(highlighted_content, unsafe_allow_html=True)
                
                # Also show clean refined version with cleaned figures
                st.subheader("Clean Refined Scope of Work")
                # Clean up figure highlighting for display
                clean_refined_sow = re.sub(r"<span style='background-color: yellow; padding: 2px 4px; border-radius: 3px;'>\[FIGURE: (.*?)\]</span>", r'<span style="background-color: yellow; padding: 2px 4px; border-radius: 3px;">\1</span>', refined_sow)
                st.markdown(clean_refined_sow, unsafe_allow_html=True)

                # Offer download for refined SoW
                docx_path = export_to_docx(refined_sow, "Refined_Scope_of_Work.docx")
                with open(docx_path, "rb") as f:
                    st.download_button("Download Refined SoW as DOCX", f, file_name="Refined_Scope_of_Work.docx")
        else:
            st.warning("Please type a suggestion before clicking 'Apply Refinement'.")
else:
    st.info("No Scope of Work to refine. Please generate one first.")
